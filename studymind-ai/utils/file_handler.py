# utils/file_handler.py
"""
File upload, validation, and text extraction.
Supports: PDF · DOCX · DOC · PPTX · PPT · TXT · MD
All file reads use UTF-8 encoding to prevent Windows errors.
"""

import os
import hashlib
from pathlib import Path
from typing import List, Dict

from dotenv import load_dotenv
load_dotenv()

UPLOAD_DIR  = Path(os.getenv("UPLOAD_DIR", "./data/uploads"))
UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
MAX_SIZE_MB = 50

# All supported extensions
SUPPORTED_EXTS = (".pdf", ".docx", ".doc", ".pptx", ".ppt", ".txt", ".md")


# ── Validation ─────────────────────────────────────────────────────────────────

def validate_file(uploaded_file) -> tuple:
    size_mb = uploaded_file.size / (1024 * 1024)
    if size_mb > MAX_SIZE_MB:
        return False, f"File too large ({size_mb:.1f} MB). Max is {MAX_SIZE_MB} MB."
    ext = Path(uploaded_file.name).suffix.lower()
    if ext not in SUPPORTED_EXTS:
        return False, (
            f"Unsupported type '{ext}'. "
            f"Allowed: PDF, DOCX, DOC, PPTX, PPT, TXT, MD."
        )
    return True, ""


# ── Save ────────────────────────────────────────────────────────────────────────

def save_uploaded_file(uploaded_file) -> Path:
    content = uploaded_file.read()
    uploaded_file.seek(0)
    h    = hashlib.md5(content).hexdigest()[:8]
    stem = Path(uploaded_file.name).stem
    ext  = Path(uploaded_file.name).suffix.lower()
    safe = f"{stem}_{h}{ext}"
    dest = UPLOAD_DIR / safe
    if not dest.exists():
        dest.write_bytes(content)
    return dest


def delete_file(file_path: str) -> bool:
    p = Path(file_path)
    if p.exists():
        p.unlink()
        return True
    return False


def list_uploaded_files() -> List[Dict]:
    files = []
    for p in UPLOAD_DIR.iterdir():
        if p.suffix.lower() in SUPPORTED_EXTS:
            files.append({
                "name":      p.name,
                "path":      str(p),
                "size_mb":   round(p.stat().st_size / (1024 * 1024), 2),
                "extension": p.suffix.lower(),
            })
    return files


# ── PDF Extraction ─────────────────────────────────────────────────────────────

def extract_text_from_pdf(file_path: str) -> Dict:
    """PyMuPDF first, pdfplumber fallback."""
    try:
        import fitz  # PyMuPDF
        doc        = fitz.open(file_path)
        pages_text = []
        for page_num, page in enumerate(doc, start=1):
            text = page.get_text("text") or ""
            pages_text.append({"page": page_num, "text": text.strip()})
        metadata = doc.metadata or {}
        doc.close()
        full_text = "\n\n".join(p["text"] for p in pages_text if p["text"])
        return {
            "text": full_text, "pages": pages_text,
            "num_pages": len(pages_text), "metadata": metadata,
            "source": file_path,
        }
    except Exception as e:
        print(f"[file_handler] PyMuPDF failed ({e}), trying pdfplumber…")
        return _pdf_plumber_fallback(file_path)


def _pdf_plumber_fallback(file_path: str) -> Dict:
    try:
        import pdfplumber
        pages_text = []
        with pdfplumber.open(file_path) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                text = page.extract_text() or ""
                pages_text.append({"page": i, "text": text.strip()})
        full_text = "\n\n".join(p["text"] for p in pages_text if p["text"])
        return {
            "text": full_text, "pages": pages_text,
            "num_pages": len(pages_text), "source": file_path,
        }
    except Exception as e:
        return {"text": f"[PDF extraction failed: {e}]", "source": file_path}


# ── DOCX / DOC Extraction ──────────────────────────────────────────────────────

def extract_text_from_docx(file_path: str) -> Dict:
    """Extract text from .docx and .doc files using python-docx."""
    try:
        import docx
        doc        = docx.Document(file_path)
        paragraphs = []
        for para in doc.paragraphs:
            t = para.text.strip()
            if t:
                paragraphs.append(t)
        # Also extract table cell text
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    t = cell.text.strip()
                    if t and t not in paragraphs:
                        paragraphs.append(t)
        full_text = "\n\n".join(paragraphs)
        return {
            "text":            full_text,
            "paragraphs":      paragraphs,
            "num_paragraphs":  len(paragraphs),
            "source":          file_path,
        }
    except Exception as e:
        return {"text": f"[DOCX/DOC extraction failed: {e}]", "source": file_path}


# ── PPTX / PPT Extraction ──────────────────────────────────────────────────────

def extract_text_from_pptx(file_path: str) -> Dict:
    """
    Extract text from PowerPoint files (.pptx, .ppt).
    Extracts text from every slide's shapes + notes.
    Requires: pip install python-pptx
    """
    try:
        from pptx import Presentation
        prs        = Presentation(file_path)
        slides_text = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts = []

            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_parts.append(shape.text.strip())

            # Extract speaker notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"[Notes] {notes_text}")

            if slide_parts:
                slides_text.append({
                    "page": slide_num,
                    "text": "\n".join(slide_parts),
                })

        full_text = "\n\n".join(s["text"] for s in slides_text if s["text"])
        return {
            "text":       full_text,
            "pages":      slides_text,
            "num_pages":  len(slides_text),
            "source":     file_path,
        }

    except ImportError:
        return {
            "text":   "[PPTX support requires python-pptx. Run: pip install python-pptx]",
            "source": file_path,
        }
    except Exception as e:
        return {"text": f"[PPTX extraction failed: {e}]", "source": file_path}


# ── TXT / MD Extraction ────────────────────────────────────────────────────────

def extract_text_from_txt(file_path: str) -> Dict:
    """UTF-8 safe text file reader — no UnicodeDecodeError on Windows."""
    text = Path(file_path).read_text(encoding="utf-8", errors="replace")
    return {"text": text, "source": file_path}


# ── Unified dispatcher ─────────────────────────────────────────────────────────

def extract_text(file_path: str) -> Dict:
    """
    Route to the correct extractor based on file extension.
    Returns a dict with at minimum: {"text": str, "source": str}
    """
    ext = Path(file_path).suffix.lower()

    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext in (".docx", ".doc"):
        return extract_text_from_docx(file_path)
    elif ext in (".pptx", ".ppt"):
        return extract_text_from_pptx(file_path)
    elif ext in (".txt", ".md"):
        return extract_text_from_txt(file_path)
    else:
        raise ValueError(
            f"Unsupported file type: {ext}. "
            f"Supported: {', '.join(SUPPORTED_EXTS)}"
        )


# ── File info helper ───────────────────────────────────────────────────────────

def get_file_info(file_path: str) -> Dict:
    """Return file metadata. Never raises — returns safe defaults on error."""
    p = Path(file_path)
    try:
        data    = extract_text(file_path)
        size_mb = round(p.stat().st_size / (1024 * 1024), 2)
        return {
            "name":       p.name,
            "path":       str(p),
            "size_mb":    size_mb,
            "extension":  p.suffix.lower(),
            "num_pages":  data.get("num_pages", "—"),
            "char_count": len(data.get("text", "")),
            "word_count": len(data.get("text", "").split()),
        }
    except Exception as e:
        return {
            "name":       p.name,
            "path":       str(p),
            "size_mb":    0,
            "extension":  p.suffix.lower(),
            "num_pages":  "—",
            "char_count": 0,
            "word_count": 0,
            "error":      str(e),
        }